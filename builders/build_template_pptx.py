"""
Build _template-path-to-zero.pptx — the 10-slide Arrive-branded template.

Output is a .pptx that, when uploaded to Google Drive, auto-converts to a
Google Slides file with full brand styling and `{{PLACEHOLDER}}` tokens
intact. The renderer later does drive.files.copy on the converted Slides
file and replaceAllText to fill placeholders.

Run:  python3 build_template_pptx.py
"""

from __future__ import annotations

import json
import pathlib

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu


SKILL_DIR = pathlib.Path(__file__).resolve().parent.parent
BRAND_DIR = SKILL_DIR / "brand"
OUT_DIR   = SKILL_DIR / "templates"
OUT_DIR.mkdir(exist_ok=True)

PALETTE = json.loads((BRAND_DIR / "palette.json").read_text())
FONT    = json.loads((BRAND_DIR / "font.json").read_text())

FONT_NAME = FONT["primary_name"]


def hex2rgb(h: str) -> RGBColor:
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


C = {k: hex2rgb(v) for k, v in PALETTE.items() if isinstance(v, str) and v.startswith("#")}


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_rect(slide, x, y, w, h, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
    shape.shadow.inherit = False
    return shape


def add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=None,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT_NAME):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color
    return tb


def add_bulleted(slide, x, y, w, h, items, *, size=12, color=None, font=FONT_NAME):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    color = color or C["dark_purple"]
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        run = p.add_run()
        run.text = f"•  {item}"
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = color
        p.space_after = Pt(6)
    return tb


def add_footer(slide, slide_num: int):
    add_rect(slide, Emu(0), SLIDE_H - Inches(0.35), SLIDE_W, Inches(0.35), C["primary_purple"])
    add_text(slide, Inches(0.4), SLIDE_H - Inches(0.33), Inches(7), Inches(0.3),
             "{{FOOTER_TEAM}}  ·  Path to Zero  ·  {{FOOTER_DATE}}",
             size=9, color=C["white"], anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, SLIDE_W - Inches(2), SLIDE_H - Inches(0.33), Inches(1.7), Inches(0.3),
             f"Confidential  ·  {slide_num} / 10",
             size=9, color=C["white"], align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def add_section_header(slide, title: str, subtitle: str | None = None):
    add_rect(slide, Emu(0), Emu(0), SLIDE_W, Inches(1.1), C["primary_purple"])
    add_rect(slide, Emu(0), Inches(1.1), SLIDE_W, Inches(0.06), C["primary_pink"])
    add_text(slide, Inches(0.4), Inches(0.18), SLIDE_W - Inches(0.8), Inches(0.55),
             title, size=26, bold=True, color=C["white"])
    if subtitle:
        add_text(slide, Inches(0.4), Inches(0.7), SLIDE_W - Inches(0.8), Inches(0.4),
                 subtitle, size=12, color=C["pink_tint_2"])
    add_text(slide, Inches(0.4), SLIDE_H - Inches(0.75), Inches(3), Inches(0.3),
             "{{TEAM_NAME}}  ·  {{AS_OF_DATE}}", size=10, color=C["muted_purple"])


# ----- slide builders -----

def build_slide_1_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_rect(slide, Emu(0), Emu(0), SLIDE_W, Inches(5.6), C["primary_purple"])
    add_rect(slide, Emu(0), Inches(4.42), SLIDE_W, Inches(0.04), C["primary_pink"])

    # Wordmark logo top-left
    slide.shapes.add_picture(str(BRAND_DIR / "arrive-wordmark.png"),
                             Inches(0.5), Inches(0.5), height=Inches(0.65))

    add_text(slide, Inches(0.5), Inches(2.0), Inches(12), Inches(1.2),
             "{{TEAM_NAME}} — Path to Zero",
             size=44, bold=True, color=C["white"])
    add_text(slide, Inches(0.5), Inches(3.1), Inches(12), Inches(0.6),
             "Critical Vulnerability Remediation",
             size=22, color=C["pink_tint_1"])
    add_text(slide, Inches(0.5), Inches(3.7), Inches(12), Inches(0.6),
             "Security Leadership Update  ·  {{AS_OF_DATE}}",
             size=16, color=C["pink_tint_2"])

    # Bottom slab
    add_rect(slide, Emu(0), Inches(5.6), SLIDE_W, SLIDE_H - Inches(5.6), C["white"])
    add_text(slide, Inches(0.5), Inches(5.85), Inches(8), Inches(0.4),
             "Engineering lead: {{ENGINEERING_LEAD}}",
             size=14, bold=True, color=C["dark_purple"])
    add_text(slide, Inches(0.5), Inches(6.3), Inches(12), Inches(0.5),
             "Q2 epics: {{EPIC_LIST}}",
             size=12, color=C["muted_purple"])

    # A-mark accent bottom-right
    slide.shapes.add_picture(str(BRAND_DIR / "arrive-amark.png"),
                             SLIDE_W - Inches(1.4), Inches(5.8), height=Inches(1.3))
    return slide


def build_slide_2_exec_summary(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "Executive Summary",
                       "{{HEADLINE_SUBTITLE}}")

    # Four stat tiles
    tile_y = Inches(1.6)
    tile_h = Inches(1.6)
    tile_w = Inches(2.9)
    gap    = Inches(0.25)
    start_x = Inches(0.4)
    labels = [
        ("{{STAT_PEAK}}",                "Peak"),
        ("{{STAT_TODAY}}",               "Today ({{AS_OF_DATE}})"),
        ("{{STAT_SERVICES_WITH_CRITS}}", "Services with criticals"),
        ("{{STAT_SERVICES_AFTER}}",      "Service after EOD"),
    ]
    for i, (val, lbl) in enumerate(labels):
        x = start_x + (tile_w + gap) * i
        add_rect(slide, x, tile_y, tile_w, tile_h, C["pink_tint_2"])
        add_text(slide, x, tile_y + Inches(0.15), tile_w, Inches(0.95), val,
                 size=46, bold=True, color=C["primary_purple"],
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, x, tile_y + Inches(1.1), tile_w, Inches(0.4), lbl,
                 size=11, color=C["muted_purple"], align=PP_ALIGN.CENTER)

    # Headline bullets
    add_text(slide, Inches(0.4), Inches(3.5), Inches(12.5), Inches(0.4),
             "Headline updates",
             size=14, bold=True, color=C["dark_purple"])
    add_bulleted(slide, Inches(0.4), Inches(3.9), Inches(12.5), Inches(3),
                 ["{{HEADLINE_BULLETS}}"], size=13)
    add_footer(slide, 2)
    return slide


def build_slide_3_burndown(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "Burndown — actual + forecast",
                       "Source: status updates and security-insights")

    # Image placeholder — a named rectangle the renderer replaces with the PNG
    ph = add_rect(slide,
                  Inches(0.4), Inches(1.5),
                  SLIDE_W - Inches(0.8), Inches(5.0),
                  C["neutral_warm"])
    ph.name = "BURNDOWN_PLACEHOLDER"
    add_text(slide, Inches(0.4), Inches(3.7), SLIDE_W - Inches(0.8), Inches(0.5),
             "{{BURNDOWN_PLACEHOLDER_CAPTION}}",
             size=12, color=C["muted_purple"], align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    add_footer(slide, 3)
    return slide


def build_slide_4_remaining(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "Remaining attack surface",
                       "{{REMAINING_SURFACE_SUBTITLE}}")

    # Table — header row + 5 placeholder rows. Renderer adds/removes rows.
    rows, cols = 6, 4
    left  = Inches(0.4)
    top   = Inches(1.7)
    width = SLIDE_W - Inches(0.8)
    height = Inches(4.5)
    tbl_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    tbl = tbl_shape.table
    tbl.columns[0].width = Inches(5.5)
    tbl.columns[1].width = Inches(1.3)
    tbl.columns[2].width = Inches(1.3)
    tbl.columns[3].width = Inches(4.4)

    headers = ["Repo", "Critical", "High", "Status"]
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = C["primary_purple"]
        tf = cell.text_frame
        tf.text = h
        r = tf.paragraphs[0].runs[0]
        r.font.name = FONT_NAME
        r.font.size = Pt(12)
        r.font.bold = True
        r.font.color.rgb = C["white"]

    for i in range(1, rows):
        bg = C["neutral_warm"] if i % 2 == 0 else C["white"]
        cells = [
            "{{ROW_%d_REPO}}" % i,
            "{{ROW_%d_CRITICAL}}" % i,
            "{{ROW_%d_HIGH}}" % i,
            "{{ROW_%d_STATUS}}" % i,
        ]
        for j, val in enumerate(cells):
            cell = tbl.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            tf = cell.text_frame
            tf.text = val
            r = tf.paragraphs[0].runs[0]
            r.font.name = FONT_NAME
            r.font.size = Pt(11)
            r.font.color.rgb = C["dark_purple"]

    add_footer(slide, 4)
    return slide


def build_slide_5_dependency(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "{{DEP_GRID_TITLE}}",
                       "Operational prerequisite status")

    add_text(slide, Inches(0.4), Inches(1.6), SLIDE_W - Inches(0.8), Inches(0.8),
             "{{DEP_GRID_NARRATIVE}}",
             size=13, color=C["dark_purple"])

    # 3-region grid (renderer adjusts column count)
    grid_y = Inches(2.7)
    grid_h = Inches(2.5)
    cols = 3
    cw = (SLIDE_W - Inches(1.0)) / cols
    for i in range(cols):
        x = Inches(0.4) + cw * i + Inches(0.1) * i
        add_rect(slide, x, grid_y, cw - Inches(0.2), grid_h, C["pink_tint_2"])
        add_text(slide, x, grid_y + Inches(0.4), cw - Inches(0.2), Inches(0.8),
                 "{{REGION_%d_NAME}}" % (i + 1),
                 size=24, bold=True, color=C["primary_purple"],
                 align=PP_ALIGN.CENTER)
        add_text(slide, x, grid_y + Inches(1.4), cw - Inches(0.2), Inches(0.8),
                 "{{REGION_%d_STATUS}}" % (i + 1),
                 size=14, color=C["dark_purple"], align=PP_ALIGN.CENTER)

    add_footer(slide, 5)
    return slide


def build_slide_6_critical_path(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "{{CP_SERVICE}} — critical path",
                       "{{CP_STAT_LINE}}")

    add_text(slide, Inches(0.4), Inches(1.6), Inches(6.0), Inches(0.4),
             "Why this is the critical path",
             size=14, bold=True, color=C["dark_purple"])
    add_bulleted(slide, Inches(0.4), Inches(2.05), Inches(8.5), Inches(4.5),
                 ["{{CP_RATIONALE_BULLETS}}"], size=13)

    # Side callout box
    add_rect(slide, Inches(9.2), Inches(1.6), Inches(3.7), Inches(5.0),
             C["pink_tint_2"])
    add_text(slide, Inches(9.4), Inches(1.8), Inches(3.3), Inches(0.5),
             "Critical surface", size=12, bold=True, color=C["primary_purple"])
    add_text(slide, Inches(9.4), Inches(2.3), Inches(3.3), Inches(1.5),
             "{{CP_STAT_LINE}}", size=14, color=C["dark_purple"])
    add_footer(slide, 6)
    return slide


def build_slide_7_epics(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "Q2 epic structure",
                       "Three workstreams")

    cols = 3
    cw = (SLIDE_W - Inches(1.0)) / cols
    top = Inches(1.7)
    h   = Inches(4.6)
    for i in range(cols):
        x = Inches(0.4) + cw * i + Inches(0.1) * i
        add_rect(slide, x, top, cw - Inches(0.2), Inches(0.8),
                 C["primary_purple"])
        add_text(slide, x, top + Inches(0.15), cw - Inches(0.2), Inches(0.5),
                 "{{EPIC_%d_ID}}" % (i + 1),
                 size=16, bold=True, color=C["white"], align=PP_ALIGN.CENTER)
        add_rect(slide, x, top + Inches(0.8), cw - Inches(0.2), h - Inches(0.8),
                 C["neutral_warm"])
        add_text(slide, x + Inches(0.15), top + Inches(1.0),
                 cw - Inches(0.5), Inches(0.5),
                 "{{EPIC_%d_STATE}}" % (i + 1),
                 size=12, bold=True, color=C["secondary_purple"])
        add_text(slide, x + Inches(0.15), top + Inches(1.5),
                 cw - Inches(0.5), h - Inches(1.7),
                 "{{EPIC_%d_SUMMARY}}" % (i + 1),
                 size=11, color=C["dark_purple"])
    add_footer(slide, 7)
    return slide


def build_slide_8_reallocation(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "Resource reallocation",
                       "Engineers reallocated from other streams")

    rows, cols = 6, 3
    left = Inches(0.4); top = Inches(1.7)
    width = SLIDE_W - Inches(0.8); height = Inches(4.5)
    tbl_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    tbl = tbl_shape.table
    tbl.columns[0].width = Inches(4.5)
    tbl.columns[1].width = Inches(4.0)
    tbl.columns[2].width = Inches(4.0)

    for j, h in enumerate(["Engineer", "From", "Allocation"]):
        cell = tbl.cell(0, j)
        cell.fill.solid(); cell.fill.fore_color.rgb = C["primary_purple"]
        tf = cell.text_frame; tf.text = h
        r = tf.paragraphs[0].runs[0]
        r.font.name = FONT_NAME; r.font.size = Pt(12); r.font.bold = True
        r.font.color.rgb = C["white"]

    for i in range(1, rows):
        bg = C["neutral_warm"] if i % 2 == 0 else C["white"]
        for j, val in enumerate([
            "{{REALLOC_%d_ENGINEER}}" % i,
            "{{REALLOC_%d_FROM}}" % i,
            "{{REALLOC_%d_TO}}" % i,
        ]):
            cell = tbl.cell(i, j)
            cell.fill.solid(); cell.fill.fore_color.rgb = bg
            tf = cell.text_frame; tf.text = val
            r = tf.paragraphs[0].runs[0]
            r.font.name = FONT_NAME; r.font.size = Pt(11)
            r.font.color.rgb = C["dark_purple"]

    add_footer(slide, 8)
    return slide


def build_slide_9_timeline(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "Updated Timeline to Zero",
                       "Dated milestones")

    rows, cols = 8, 3
    left = Inches(0.4); top = Inches(1.7)
    width = SLIDE_W - Inches(0.8); height = Inches(5.0)
    tbl_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    tbl = tbl_shape.table
    tbl.columns[0].width = Inches(1.8)
    tbl.columns[1].width = Inches(8.7)
    tbl.columns[2].width = Inches(2.0)

    for j, h in enumerate(["Date", "Milestone", "Open criticals"]):
        cell = tbl.cell(0, j)
        cell.fill.solid(); cell.fill.fore_color.rgb = C["primary_purple"]
        tf = cell.text_frame; tf.text = h
        r = tf.paragraphs[0].runs[0]
        r.font.name = FONT_NAME; r.font.size = Pt(12); r.font.bold = True
        r.font.color.rgb = C["white"]

    for i in range(1, rows):
        bg = C["neutral_warm"] if i % 2 == 0 else C["white"]
        for j, val in enumerate([
            "{{TIMELINE_%d_DATE}}" % i,
            "{{TIMELINE_%d_MILESTONE}}" % i,
            "{{TIMELINE_%d_OPEN}}" % i,
        ]):
            cell = tbl.cell(i, j)
            cell.fill.solid(); cell.fill.fore_color.rgb = bg
            tf = cell.text_frame; tf.text = val
            r = tf.paragraphs[0].runs[0]
            r.font.name = FONT_NAME; r.font.size = Pt(11)
            r.font.color.rgb = C["dark_purple"]

    add_footer(slide, 9)
    return slide


def build_slide_10_asks(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "What we need from Security leadership",
                       "Asks")

    top = Inches(1.7)
    h = Inches(1.1)
    gap = Inches(0.2)
    for i in range(4):
        y = top + (h + gap) * i
        # Number circle
        add_rect(slide, Inches(0.4), y, Inches(1.0), h, C["primary_purple"])
        add_text(slide, Inches(0.4), y, Inches(1.0), h,
                 "{{ASK_%d_NUM}}" % (i + 1),
                 size=32, bold=True, color=C["white"],
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Body card
        add_rect(slide, Inches(1.5), y, SLIDE_W - Inches(1.9), h, C["neutral_warm"])
        add_text(slide, Inches(1.7), y + Inches(0.1), SLIDE_W - Inches(2.1), Inches(0.4),
                 "{{ASK_%d_TITLE}}" % (i + 1),
                 size=14, bold=True, color=C["primary_purple"])
        add_text(slide, Inches(1.7), y + Inches(0.5), SLIDE_W - Inches(2.1), Inches(0.6),
                 "{{ASK_%d_BODY}}" % (i + 1),
                 size=11, color=C["dark_purple"])

    add_footer(slide, 10)
    return slide


def main():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    build_slide_1_title(prs)
    build_slide_2_exec_summary(prs)
    build_slide_3_burndown(prs)
    build_slide_4_remaining(prs)
    build_slide_5_dependency(prs)
    build_slide_6_critical_path(prs)
    build_slide_7_epics(prs)
    build_slide_8_reallocation(prs)
    build_slide_9_timeline(prs)
    build_slide_10_asks(prs)

    out = OUT_DIR / "_template-path-to-zero.pptx"
    prs.save(str(out))
    print(f"wrote {out}  ({out.stat().st_size:,} bytes)")


if __name__ == "__main__":
    main()
