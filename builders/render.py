"""
Render a path-to-zero deliverable pair from a plan.json.

Fills the PPTX and DOCX templates with values from plan.json, replacing
`{{TOKEN}}` placeholders inline (preserving formatting), inserting the
burndown PNG, and adjusting tables to match data row counts.

Outputs land in --out-dir (default: <cwd>/path-to-zero-plan/).
  <team>-Path-to-Zero-<YYYY-MM-DD>.pptx
  <team>-Security-Leadership-Update-<YYYY-MM-DD>.docx

Run:
  python3 render.py path/to/plan.json
  python3 render.py path/to/plan.json --out-dir /tmp/preview
"""

from __future__ import annotations

import argparse
import copy
import json
import pathlib
import re
import sys

from pptx import Presentation
from pptx.util import Emu
from docx import Document


SKILL_DIR = pathlib.Path(__file__).resolve().parent.parent
TEMPLATES = SKILL_DIR / "templates"
SCHEMAS   = SKILL_DIR / "schemas"

PPTX_TEMPLATE = TEMPLATES / "_template-path-to-zero.pptx"
DOCX_TEMPLATE = TEMPLATES / "_template-leadership-update.docx"

SCHEMA_VERSION_SUPPORTED = 1


# ----- plan.json loading + minimal validation -----

REQUIRED_KEYS = {
    "schema_version", "team", "as_of_date", "engineering_lead", "epics",
    "exec_summary", "burndown_png_path", "remaining_surface",
    "dependency_grid", "critical_path_deep_dive", "epic_status",
    "resource_reallocation", "timeline", "asks",
}


def load_plan(path: pathlib.Path) -> dict:
    plan = json.loads(path.read_text())
    missing = REQUIRED_KEYS - plan.keys()
    if missing:
        raise SystemExit(f"plan.json missing required keys: {sorted(missing)}")
    if plan["schema_version"] != SCHEMA_VERSION_SUPPORTED:
        raise SystemExit(
            f"plan.json schema_version={plan['schema_version']} not supported "
            f"(this renderer supports v{SCHEMA_VERSION_SUPPORTED})"
        )
    return plan


# ----- build the token map from plan.json -----

def build_token_map(plan: dict) -> dict[str, str]:
    t: dict[str, str] = {}

    t["TEAM_NAME"]          = plan["team"]
    t["AS_OF_DATE"]         = plan["as_of_date"]
    t["ENGINEERING_LEAD"]   = plan["engineering_lead"]
    t["FOOTER_TEAM"]        = plan["team"]
    t["FOOTER_DATE"]        = plan["as_of_date"]
    t["EPIC_LIST"]          = "  ·  ".join(
        f"{e['id']} ({e['name']})" for e in plan["epics"]
    )

    es = plan["exec_summary"]
    t["STAT_PEAK"]                = str(es["peak"])
    t["STAT_TODAY"]               = str(es["today"])
    t["STAT_SERVICES_WITH_CRITS"] = str(es["services_with_criticals"])
    t["STAT_SERVICES_AFTER"]      = str(es["services_after_eod"])
    t["HEADLINE_SUBTITLE"]        = es.get("subtitle", "")
    t["HEADLINE_BULLETS"]         = "\n".join(f"• {b}" for b in es["headline_bullets"])

    t["BURNDOWN_PLACEHOLDER_CAPTION"] = (
        f"Burndown: {plan['team']} criticals through {plan['as_of_date']}"
    )

    t["REMAINING_SURFACE_SUBTITLE"] = (
        f"{len(plan['remaining_surface'])} repo(s) with open criticals"
    )
    # Pre-allocated 5 surface rows in the PPTX, same in the DOCX
    for i in range(1, 6):
        if i <= len(plan["remaining_surface"]):
            row = plan["remaining_surface"][i - 1]
            t[f"ROW_{i}_REPO"]     = row["repo"]
            t[f"ROW_{i}_CRITICAL"] = str(row["critical"])
            t[f"ROW_{i}_HIGH"]     = str(row["high"])
            t[f"ROW_{i}_STATUS"]   = row["status"]
        else:
            t[f"ROW_{i}_REPO"] = t[f"ROW_{i}_CRITICAL"] = ""
            t[f"ROW_{i}_HIGH"] = t[f"ROW_{i}_STATUS"] = ""

    dg = plan["dependency_grid"]
    t["DEP_GRID_TITLE"]     = dg["title"]
    t["DEP_GRID_NARRATIVE"] = dg["narrative"]
    for i in range(1, 4):
        if i <= len(dg["regions"]):
            r = dg["regions"][i - 1]
            t[f"REGION_{i}_NAME"]   = r["name"]
            t[f"REGION_{i}_STATUS"] = r["status"]
        else:
            t[f"REGION_{i}_NAME"] = t[f"REGION_{i}_STATUS"] = ""

    cp = plan["critical_path_deep_dive"]
    t["CP_SERVICE"]           = cp["service"]
    t["CP_STAT_LINE"]         = cp["stat_line"]
    t["CP_RATIONALE_BULLETS"] = "\n".join(f"• {b}" for b in cp["rationale_bullets"])

    for i in range(1, 4):
        if i <= len(plan["epic_status"]):
            e = plan["epic_status"][i - 1]
            t[f"EPIC_{i}_ID"]      = e["id"]
            t[f"EPIC_{i}_STATE"]   = e["state"]
            t[f"EPIC_{i}_SUMMARY"] = e["summary"]
        else:
            t[f"EPIC_{i}_ID"] = t[f"EPIC_{i}_STATE"] = t[f"EPIC_{i}_SUMMARY"] = ""

    for i in range(1, 6):
        if i <= len(plan["resource_reallocation"]):
            r = plan["resource_reallocation"][i - 1]
            t[f"REALLOC_{i}_ENGINEER"] = r["engineer"]
            t[f"REALLOC_{i}_FROM"]     = r["from"]
            t[f"REALLOC_{i}_TO"]       = r["to"]
        else:
            t[f"REALLOC_{i}_ENGINEER"] = ""
            t[f"REALLOC_{i}_FROM"]     = ""
            t[f"REALLOC_{i}_TO"]       = ""

    for i in range(1, 8):
        if i <= len(plan["timeline"]):
            r = plan["timeline"][i - 1]
            t[f"TIMELINE_{i}_DATE"]      = r["date"]
            t[f"TIMELINE_{i}_MILESTONE"] = r["milestone"]
            t[f"TIMELINE_{i}_OPEN"]      = str(r["open_criticals"])
        else:
            t[f"TIMELINE_{i}_DATE"] = ""
            t[f"TIMELINE_{i}_MILESTONE"] = ""
            t[f"TIMELINE_{i}_OPEN"] = ""

    for i in range(1, 5):
        if i <= len(plan["asks"]):
            a = plan["asks"][i - 1]
            t[f"ASK_{i}_NUM"]   = str(a["num"])
            t[f"ASK_{i}_TITLE"] = a["title"]
            t[f"ASK_{i}_BODY"]  = a["body"]
        else:
            t[f"ASK_{i}_NUM"] = t[f"ASK_{i}_TITLE"] = t[f"ASK_{i}_BODY"] = ""

    return t


# ----- substitution helpers -----

TOKEN_RE = re.compile(r"\{\{([A-Z_0-9]+)\}\}")


def sub(text: str, tokens: dict[str, str]) -> str:
    def replace(m):
        key = m.group(1)
        return tokens.get(key, m.group(0))
    return TOKEN_RE.sub(replace, text)


def fill_runs(runs, tokens: dict[str, str]):
    """Replace tokens across a sequence of runs that share a paragraph.

    Tokens may span runs after PowerPoint/Word rewrites. We collapse the
    paragraph text, substitute, then write back into the first run and
    blank the rest. This loses intra-paragraph mixed formatting but the
    builder uses uniform formatting per paragraph, so it's safe.
    """
    if not runs:
        return
    joined = "".join(r.text for r in runs)
    if "{{" not in joined:
        return
    replaced = sub(joined, tokens)
    runs[0].text = replaced
    for r in runs[1:]:
        r.text = ""


# ----- PPTX rendering -----

def render_pptx(plan: dict, tokens: dict[str, str], out_path: pathlib.Path):
    prs = Presentation(str(PPTX_TEMPLATE))

    burndown_png = pathlib.Path(plan["burndown_png_path"])

    for slide in prs.slides:
        # Replace placeholder rectangle with burndown image
        for shape in list(slide.shapes):
            if shape.name == "BURNDOWN_PLACEHOLDER" and burndown_png.exists():
                left, top, width, height = shape.left, shape.top, shape.width, shape.height
                sp = shape._element
                sp.getparent().remove(sp)
                slide.shapes.add_picture(str(burndown_png), left, top, width, height)

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    fill_runs(para.runs, tokens)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for para in cell.text_frame.paragraphs:
                            fill_runs(para.runs, tokens)

    prs.save(str(out_path))


# ----- DOCX rendering -----

def render_docx(plan: dict, tokens: dict[str, str], out_path: pathlib.Path):
    doc = Document(str(DOCX_TEMPLATE))

    burndown_png = pathlib.Path(plan["burndown_png_path"])

    # Paragraphs at body level
    for para in doc.paragraphs:
        if "{{BURNDOWN_PLACEHOLDER}}" in para.text and burndown_png.exists():
            for r in para.runs:
                r.text = ""
            para.runs[0] if para.runs else para.add_run()
            (para.runs[0] if para.runs else para.add_run()).add_picture(
                str(burndown_png), width=Emu(5_500_000)
            )
            continue
        fill_runs(para.runs, tokens)

    # Tables
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    fill_runs(para.runs, tokens)

    doc.save(str(out_path))


# ----- entrypoint -----

def main(argv):
    ap = argparse.ArgumentParser()
    ap.add_argument("plan", help="path to plan.json")
    ap.add_argument("--out-dir", default=None,
                    help="output dir (default: <cwd>/path-to-zero-plan)")
    args = ap.parse_args(argv)

    plan_path = pathlib.Path(args.plan).resolve()
    plan = load_plan(plan_path)
    tokens = build_token_map(plan)

    out_dir = pathlib.Path(args.out_dir) if args.out_dir else pathlib.Path.cwd() / "path-to-zero-plan"
    out_dir.mkdir(parents=True, exist_ok=True)

    team = plan["team"]
    date = plan["as_of_date"]
    pptx_out = out_dir / f"{team}-Path-to-Zero-{date}.pptx"
    docx_out = out_dir / f"{team}-Security-Leadership-Update-{date}.docx"

    render_pptx(plan, tokens, pptx_out)
    print(f"wrote {pptx_out}  ({pptx_out.stat().st_size:,} bytes)")

    render_docx(plan, tokens, docx_out)
    print(f"wrote {docx_out}  ({docx_out.stat().st_size:,} bytes)")


if __name__ == "__main__":
    main(sys.argv[1:])
